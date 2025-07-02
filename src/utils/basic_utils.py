import os
import json
import langid
import random
import aiohttp
import json_repair
import pandas as pd
import PyPDF2
from pptx import Presentation


def get_clear_model_name(model_name):
    if "qwq" in model_name.lower():
        return "qwq"
    elif "dpsk" in model_name.lower():
        return "dpsk"
    elif "llama" in model_name.lower():
        return "llama"
    else:
        return model_name


def load_data(dataset_name, split, subset_num=-1):
    data_path = f"../data/{dataset_name}/{split}.json"

    print("-----------------------")
    print(f"Using {dataset_name} {split} set.")
    print("-----------------------")

    # Load and prepare data
    with open(data_path, "r", encoding="utf-8") as json_file:
        data = json.load(json_file)

    if subset_num != -1:
        indices = list(range(len(data)))
        selected_indices = random.sample(indices, min(subset_num, len(indices)))
        data = [data[i] for i in selected_indices]

    return data


def prepare_output_dir(model_name, dataset_name, experiment_name=None, use_single_dir=False, save_dir="./outputs"):
    # Define output directory
    if "qwq" in model_name.lower():
        model_short_name = "qwq"
    elif "deepseek" in model_name.lower():
        if "llama-8b" in model_name.lower():
            model_short_name = "dpsk-llama-8b"
        elif "llama-70b" in model_name.lower():
            model_short_name = "dpsk-llama-70b"
        elif "qwen-1.5b" in model_name.lower():
            model_short_name = "dpsk-qwen-1.5b"
        elif "qwen-7b" in model_name.lower():
            model_short_name = "dpsk-qwen-7b"
        elif "qwen-32b" in model_name.lower():
            model_short_name = "dpsk-qwen-32b"
    elif "sky-t1" in model_name.lower():
        model_short_name = "sky-t1"
    else:
        model_short_name = model_name.split("/")[-1].lower().replace("-instruct", "")
    if experiment_name is None:
        experiment_name = "webthinker"
    else:
        experiment_name = str(experiment_name)
    if use_single_dir:
        # put the series of experiments in the same directory
        output_dir = f"{save_dir}/{experiment_name}/{dataset_name}.{model_short_name}"
    else:
        output_dir = f"{save_dir}/{dataset_name}.{model_short_name}.{experiment_name}"
    os.makedirs(output_dir, exist_ok=True)
    return output_dir


async def load_lora_adapter(api_base_url: str, lora_name: str, lora_path: str) -> bool:
    try:
        lora_load_url = f"{api_base_url}/load_lora_adapter"
        lora_payload = {"lora_name": lora_name, "lora_path": lora_path}
        async with aiohttp.ClientSession() as session:
            async with session.post(lora_load_url, json=lora_payload) as response:
                return response.status == 200
    except Exception as e:
        print(f"Error loading LoRA adapter: {e}")
        return False


async def unload_lora_adapter(api_base_url: str, lora_name: str) -> bool:
    try:
        unload_url = f"{api_base_url}/unload_lora_adapter"
        unload_payload = {"lora_name": lora_name}
        async with aiohttp.ClientSession() as session:
            async with session.post(unload_url, json=unload_payload) as response:
                return response.status == 200
    except Exception as e:
        print(f"Error unloading LoRA adapter: {e}")
        return False


async def check_lora_and_load(lora_name, lora_path, api_base_url):
    if lora_name and lora_path:
        print(f"Loading LoRA adapter '{lora_name}' from {lora_path}")
        success = await load_lora_adapter(api_base_url, lora_name, lora_path)
        if not success:
            print("Failed to load LoRA adapter")
            return
        else:
            print("LoRA adapter loaded successfully")


async def check_lora_and_unload(lora_name, api_base_url):
    if lora_name:
        print(f"Unloading LoRA adapter '{lora_name}'")
        await unload_lora_adapter(api_base_url, lora_name)


def load_caches(cache_dir, keep_links):
    search_cache_path = os.path.join(cache_dir, "search_cache.json")
    if keep_links:
        url_cache_path = os.path.join(cache_dir, "url_cache_with_links.json")
    else:
        url_cache_path = os.path.join(cache_dir, "url_cache.json")
    file_cache_path = os.path.join(cache_dir, "file_cache.json")

    os.makedirs(cache_dir, exist_ok=True)

    # Load existing caches
    search_cache = (
        json.load(open(search_cache_path, "r", encoding="utf-8")) if os.path.exists(search_cache_path) else {}
    )
    url_cache = json.load(open(url_cache_path)) if os.path.exists(url_cache_path) else {}
    file_cache = json.load(open(file_cache_path)) if os.path.exists(file_cache_path) else {}
    return search_cache, url_cache, file_cache


def save_caches(cache_dir, search_cache, url_cache, file_cache):
    search_cache_path = os.path.join(cache_dir, "search_cache.json")
    url_cache_path = os.path.join(cache_dir, "url_cache_with_links.json")
    file_cache_path = os.path.join(cache_dir, "file_cache.json")

    with open(search_cache_path, "w", encoding="utf-8") as f:
        json.dump(search_cache, f, ensure_ascii=False, indent=2)
    with open(url_cache_path, "w", encoding="utf-8") as f:
        json.dump(url_cache, f, ensure_ascii=False, indent=2)
    with open(file_cache_path, "w", encoding="utf-8") as f:
        json.dump(file_cache, f, ensure_ascii=False, indent=2)


def get_file_content(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    try:
        if file_extension in [".txt", ".py"]:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif file_extension == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string()

        elif file_extension == ".xlsx":
            df = pd.read_excel(file_path)
            return df.to_string()

        elif file_extension == ".pptx":
            text_content = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return "\n".join(text_content)

        elif file_extension == ".pdf":
            text_content = []
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
            return "\n".join(text_content)

        else:
            return f"No valid file content"

    except Exception as e:
        return ""


def judge_zh(input_str: str):
    assert isinstance(input_str, str), input_str
    if len(input_str) == 0:
        return False
    detect_result = langid.classify(input_str)
    if detect_result[0] == "zh":
        return True
    else:
        return False
